from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PptxRGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from io import BytesIO
from typing import Dict, Any, List, Tuple
import json
import random


class ExportService:
    """
    Servicio para exportar actividades a Word, Excel y PowerPoint con diseños profesionales
    """

    @staticmethod
    def _get_excel_column_letter(col_idx):
        """Convierte índice de columna (0-based) a letra de Excel (A, B, ..., Z, AA, AB, ...)"""
        result = ""
        while col_idx >= 0:
            result = chr(65 + (col_idx % 26)) + result
            col_idx = col_idx // 26 - 1
        return result

    # Esquemas de colores profesionales tipo Gamma AI
    COLOR_SCHEMES = {
        "modern_blue": {
            "primary": (54, 96, 146),      # Azul profesional
            "secondary": (91, 155, 213),    # Azul claro
            "accent": (255, 192, 0),        # Amarillo dorado
            "background": (247, 249, 252),  # Gris muy claro
            "text_dark": (31, 41, 55),      # Gris oscuro
            "text_light": (107, 114, 128),  # Gris medio
        },
        "elegant_purple": {
            "primary": (109, 40, 217),      # Morado vibrante
            "secondary": (167, 139, 250),   # Morado claro
            "accent": (251, 146, 60),       # Naranja
            "background": (249, 250, 251),  # Blanco humo
            "text_dark": (17, 24, 39),      # Casi negro
            "text_light": (75, 85, 99),     # Gris
        },
        "fresh_green": {
            "primary": (16, 185, 129),      # Verde menta
            "secondary": (110, 231, 183),   # Verde claro
            "accent": (245, 158, 11),       # Amarillo mostaza
            "background": (240, 253, 250),  # Verde muy claro
            "text_dark": (6, 78, 59),       # Verde oscuro
            "text_light": (52, 211, 153),   # Verde medio
        },
        "warm_orange": {
            "primary": (249, 115, 22),      # Naranja
            "secondary": (251, 146, 60),    # Naranja claro
            "accent": (234, 88, 12),        # Naranja oscuro
            "background": (255, 251, 235),  # Crema
            "text_dark": (124, 45, 18),     # Marrón
            "text_light": (161, 98, 7),     # Marrón claro
        },
        "corporate_gray": {
            "primary": (71, 85, 105),       # Gris azulado
            "secondary": (148, 163, 184),   # Gris claro
            "accent": (59, 130, 246),       # Azul brillante
            "background": (248, 250, 252),  # Casi blanco
            "text_dark": (15, 23, 42),      # Gris muy oscuro
            "text_light": (100, 116, 139),  # Gris medio
        },
    }

    def export_to_word(self, activity_data: Dict[str, Any], activity_type: str) -> BytesIO:
        """
        Exporta una actividad a formato Word
        """
        document = Document()

        # Título
        title = document.add_heading(activity_data.get("title", "Actividad Educativa"), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Descripción
        if "description" in activity_data:
            document.add_paragraph(activity_data["description"])
            document.add_paragraph()

        # Contenido específico según el tipo
        if activity_type == "exam":
            self._add_exam_to_word(document, activity_data)
        elif activity_type == "summary":
            self._add_summary_to_word(document, activity_data)
        elif activity_type == "class_activity":
            self._add_class_activity_to_word(document, activity_data)
        elif activity_type == "rubric":
            self._add_rubric_to_word(document, activity_data)
        elif activity_type == "writing_correction":
            self._add_correction_to_word(document, activity_data)
        elif activity_type == "slides":
            self._add_slides_to_word(document, activity_data)
        elif activity_type == "email":
            self._add_email_to_word(document, activity_data)
        elif activity_type == "survey":
            self._add_survey_to_word(document, activity_data)
        elif activity_type == "story":
            self._add_story_to_word(document, activity_data)
        elif activity_type == "crossword":
            self._add_crossword_to_word(document, activity_data)
        elif activity_type == "word_search":
            self._add_word_search_to_word(document, activity_data)

        # Guardar en buffer
        buffer = BytesIO()
        document.save(buffer)
        buffer.seek(0)
        return buffer

    def export_to_excel(self, activity_data: Dict[str, Any], activity_type: str) -> BytesIO:
        """
        Exporta una actividad a formato Excel
        """
        wb = Workbook()
        ws = wb.active

        # Estilo para encabezados
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        # Título
        ws['A1'] = activity_data.get("title", "Actividad Educativa")
        ws['A1'].font = Font(size=16, bold=True)
        ws.merge_cells('A1:D1')

        # Contenido específico según el tipo
        if activity_type == "exam":
            self._add_exam_to_excel(ws, activity_data, header_fill, header_font)
        elif activity_type == "survey":
            self._add_survey_to_excel(ws, activity_data, header_fill, header_font)
        elif activity_type == "rubric":
            self._add_rubric_to_excel(ws, activity_data, header_fill, header_font)
        elif activity_type == "crossword":
            self._add_crossword_to_excel(ws, activity_data, header_fill, header_font)
        elif activity_type == "word_search":
            self._add_word_search_to_excel(ws, activity_data, header_fill, header_font)

        # Ajustar ancho de columnas (evitar MergedCell)
        from openpyxl.cell.cell import MergedCell
        for column_cells in ws.columns:
            max_length = 0
            column_letter = None

            for cell in column_cells:
                # Saltar celdas combinadas
                if isinstance(cell, MergedCell):
                    continue

                # Obtener la letra de la columna
                if column_letter is None:
                    column_letter = cell.column_letter

                try:
                    cell_value = str(cell.value) if cell.value is not None else ""
                    if len(cell_value) > max_length:
                        max_length = len(cell_value)
                except:
                    pass

            # Ajustar ancho solo si encontramos una columna válida
            if column_letter:
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

        # Guardar en buffer
        buffer = BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer

    def export_to_pptx(self, activity_data: Dict[str, Any], activity_type: str) -> BytesIO:
        """
        Exporta una actividad de tipo 'slides' a formato PowerPoint (.pptx) con diseño profesional estilo Gamma AI
        """
        if activity_type != "slides":
            raise ValueError("Sólo se puede exportar a PPTX actividades de tipo 'slides'")

        prs = Presentation()

        # Configurar tamaño de diapositiva (16:9)
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(5.625)

        try:
            content = json.loads(activity_data.get("content", "{}")) if isinstance(activity_data.get("content"), str) else activity_data.get("content", {})
        except:
            content = {}

        # Seleccionar esquema de color aleatorio
        scheme_name = random.choice(list(self.COLOR_SCHEMES.keys()))
        colors = self.COLOR_SCHEMES[scheme_name]

        slides_data = content.get("slides", [])

        for idx, slide_data in enumerate(slides_data):
            is_title_slide = idx == 0

            # Usar layout en blanco para diseño personalizado
            layout = prs.slide_layouts[6]  # Layout en blanco
            slide = prs.slides.add_slide(layout)

            if is_title_slide:
                self._create_title_slide(slide, slide_data, colors, activity_data.get("title", ""))
            else:
                self._create_content_slide(slide, slide_data, colors)

            # Notas de la diapositiva
            if "notes" in slide_data:
                try:
                    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
                except:
                    pass

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def _create_title_slide(self, slide, slide_data: Dict, colors: Dict, main_title: str):
        """
        Crea la diapositiva de título con diseño profesional
        """
        # Fondo de color sólido
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PptxRGBColor(*colors["primary"])

        # Barra decorativa superior
        left = PptxInches(0)
        top = PptxInches(0)
        width = PptxInches(10)
        height = PptxInches(0.5)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(*colors["accent"])
        shape.line.fill.background()

        # Título principal (centrado verticalmente)
        left = PptxInches(1)
        top = PptxInches(2)
        width = PptxInches(8)
        height = PptxInches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = slide_data.get("title", main_title)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = PptxPt(54)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(255, 255, 255)
        p.font.name = "Segoe UI"

        # Subtítulo o descripción
        if slide_data.get("content"):
            left = PptxInches(1.5)
            top = PptxInches(3.8)
            width = PptxInches(7)
            height = PptxInches(0.8)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = subtitle_box.text_frame

            content_items = slide_data.get("content", [])
            if content_items:
                p = text_frame.paragraphs[0]
                p.text = content_items[0] if len(content_items) > 0 else ""
                p.alignment = PP_ALIGN.CENTER
                p.font.size = PptxPt(20)
                p.font.color.rgb = PptxRGBColor(*colors["background"])
                p.font.name = "Segoe UI Light"

    def _create_content_slide(self, slide, slide_data: Dict, colors: Dict):
        """
        Crea una diapositiva de contenido con diseño moderno y limpio
        """
        # Fondo claro
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PptxRGBColor(255, 255, 255)

        # Barra lateral decorativa (acento)
        left = PptxInches(0)
        top = PptxInches(0)
        width = PptxInches(0.15)
        height = PptxInches(5.625)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(*colors["accent"])
        shape.line.fill.background()

        # Título de la diapositiva
        left = PptxInches(0.5)
        top = PptxInches(0.4)
        width = PptxInches(9)
        height = PptxInches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(*colors["primary"])
        p.font.name = "Segoe UI"

        # Línea decorativa bajo el título
        left = PptxInches(0.5)
        top = PptxInches(1.3)
        width = PptxInches(2)
        height = PptxInches(0.05)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(*colors["secondary"])
        shape.line.fill.background()

        # Contenido (puntos clave con diseño tipo card)
        content_items = slide_data.get("content", [])

        # Determinar layout basado en cantidad de items
        if len(content_items) <= 3:
            # Layout vertical con cards grandes
            self._add_vertical_content(slide, content_items, colors)
        else:
            # Layout de lista con bullets modernos
            self._add_list_content(slide, content_items, colors)

    def _add_vertical_content(self, slide, items: List[str], colors: Dict):
        """
        Agrega contenido en formato vertical con cards tipo Gamma
        """
        card_height = PptxInches(1.2)
        card_spacing = PptxInches(0.25)
        start_top = PptxInches(1.8)

        for idx, item in enumerate(items[:3]):  # Máximo 3 items
            # Card con fondo suave
            left = PptxInches(1)
            top = start_top + (card_height + card_spacing) * idx
            width = PptxInches(8)

            # Fondo del card
            card = slide.shapes.add_shape(1, left, top, width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = PptxRGBColor(*colors["background"])
            card.line.color.rgb = PptxRGBColor(*colors["secondary"])
            card.line.width = PptxPt(1)

            # Número o bullet decorativo
            bullet_left = PptxInches(1.2)
            bullet_top = top + PptxInches(0.15)
            bullet_size = PptxInches(0.5)

            bullet = slide.shapes.add_shape(1, bullet_left, bullet_top, bullet_size, bullet_size)
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = PptxRGBColor(*colors["accent"])
            bullet.line.fill.background()

            # Texto del número
            text_frame = bullet.text_frame
            p = text_frame.paragraphs[0]
            p.text = str(idx + 1)
            p.alignment = PP_ALIGN.CENTER
            p.font.size = PptxPt(20)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(255, 255, 255)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Texto del contenido
            text_left = PptxInches(2)
            text_top = top + PptxInches(0.2)
            text_width = PptxInches(6.5)
            text_height = PptxInches(0.8)

            text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.text = item
            p.font.size = PptxPt(18)
            p.font.color.rgb = PptxRGBColor(*colors["text_dark"])
            p.font.name = "Segoe UI"
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_list_content(self, slide, items: List[str], colors: Dict):
        """
        Agrega contenido en formato lista con bullets modernos
        """
        left = PptxInches(1)
        top = PptxInches(1.8)
        width = PptxInches(8)
        height = PptxInches(3.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for idx, item in enumerate(items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = "  " + item  # Espaciado para el bullet custom
            p.level = 0
            p.font.size = PptxPt(16)
            p.font.color.rgb = PptxRGBColor(*colors["text_dark"])
            p.font.name = "Segoe UI"
            p.space_before = PptxPt(12)

            # Bullet personalizado (círculo con color de acento)
            p.font.bullet_char = "●"
            p.font.bullet_color_rgb = PptxRGBColor(*colors["accent"])

    # Métodos auxiliares para Word
    def _add_exam_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        if "instructions" in content:
            doc.add_heading("Instrucciones", 2)
            doc.add_paragraph(content["instructions"])

        doc.add_heading("Preguntas", 2)
        for q in content.get("questions", []):
            p = doc.add_paragraph()
            p.add_run(f"Pregunta {q.get('id', 0)}: ").bold = True
            p.add_run(q.get("question", ""))

            if q.get("type") == "multiple_choice" and "options" in q:
                for opt in q["options"]:
                    doc.add_paragraph(f"  {opt}", style='List Bullet')

            doc.add_paragraph()

    def _add_summary_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        doc.add_heading("Resumen", 2)
        doc.add_paragraph(content.get("summary", ""))

        if "key_points" in content:
            doc.add_heading("Puntos Clave", 2)
            for point in content["key_points"]:
                doc.add_paragraph(point, style='List Bullet')

    def _add_class_activity_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        if "objectives" in content:
            doc.add_heading("Objetivos", 2)
            for obj in content["objectives"]:
                doc.add_paragraph(obj, style='List Bullet')

        if "materials" in content:
            doc.add_heading("Materiales", 2)
            for mat in content["materials"]:
                doc.add_paragraph(mat, style='List Bullet')

        if "steps" in content:
            doc.add_heading("Pasos de la Actividad", 2)
            for step in content["steps"]:
                p = doc.add_paragraph()
                p.add_run(f"Paso {step.get('step', 0)} ({step.get('duration_minutes', 0)} min): ").bold = True
                p.add_run(step.get("description", ""))

    def _add_rubric_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        doc.add_heading("Rúbrica de Evaluación", 2)

        for criterion in content.get("criteria", []):
            doc.add_heading(f"{criterion.get('name', '')} ({criterion.get('weight', 0)}%)", 3)

            table = doc.add_table(rows=1, cols=3)
            table.style = 'Light Grid Accent 1'

            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Nivel'
            hdr_cells[1].text = 'Puntos'
            hdr_cells[2].text = 'Descripción'

            for level in criterion.get("levels", []):
                row_cells = table.add_row().cells
                row_cells[0].text = level.get("level", "")
                row_cells[1].text = str(level.get("points", ""))
                row_cells[2].text = level.get("description", "")

            doc.add_paragraph()

    def _add_correction_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        doc.add_heading("Texto Original", 2)
        doc.add_paragraph(content.get("original_text", ""))

        doc.add_heading("Texto Corregido", 2)
        doc.add_paragraph(content.get("corrected_text", ""))

        if "errors" in content and content["errors"]:
            doc.add_heading("Errores Detectados", 2)
            for error in content["errors"]:
                p = doc.add_paragraph()
                p.add_run(f"{error.get('type', '').capitalize()}: ").bold = True
                p.add_run(f"'{error.get('original', '')}' → '{error.get('correction', '')}'\n")
                p.add_run(f"Explicación: {error.get('explanation', '')}")

    def _add_slides_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        for slide in content.get("slides", []):
            doc.add_heading(f"Diapositiva {slide.get('slide_number', 0)}: {slide.get('title', '')}", 2)

            for point in slide.get("content", []):
                doc.add_paragraph(point, style='List Bullet')

            if "notes" in slide:
                p = doc.add_paragraph()
                p.add_run("Notas: ").italic = True
                p.add_run(slide["notes"])

            doc.add_paragraph()

    def _add_email_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        doc.add_heading("Asunto", 2)
        doc.add_paragraph(content.get("subject", ""))

        doc.add_heading("Cuerpo del Mensaje", 2)
        doc.add_paragraph(content.get("body", ""))

        doc.add_paragraph()
        doc.add_paragraph(content.get("closing", ""))

    def _add_survey_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        if "description" in content:
            doc.add_paragraph(content["description"])

        doc.add_heading("Preguntas", 2)
        for q in content.get("questions", []):
            p = doc.add_paragraph()
            p.add_run(f"{q.get('id', 0)}. {q.get('question', '')}").bold = True

            if q.get("type") == "multiple_choice" and "options" in q:
                for opt in q["options"]:
                    doc.add_paragraph(f"  {opt}", style='List Bullet')

            doc.add_paragraph()

    def _add_story_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        doc.add_paragraph(content.get("story", ""))

        if "moral" in content:
            doc.add_heading("Moraleja", 2)
            doc.add_paragraph(content["moral"])

    def _add_crossword_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except Exception as e:
            print(f"Error parseando contenido del crucigrama para Word: {e}")
            content = {}

        print(f"Contenido del crucigrama para Word: {content}")

        # Descripción si existe
        if "description" in content:
            doc.add_paragraph(content["description"])
            doc.add_paragraph()

        clues = content.get("clues", {})

        if not clues or (not clues.get("across", []) and not clues.get("down", [])):
            doc.add_heading("Error", 2)
            doc.add_paragraph("No se encontraron pistas para este crucigrama.")
            return

        across_clues = clues.get("across", [])
        down_clues = clues.get("down", [])

        # Generar la cuadrícula visual del crucigrama
        grid_size = content.get("grid_size", {"rows": 15, "cols": 15})
        rows = grid_size.get("rows", 15)
        cols = grid_size.get("cols", 15)

        # Crear matriz de la cuadrícula
        grid = [[None for _ in range(cols)] for _ in range(rows)]

        # Marcar las posiciones de las palabras horizontales
        for clue in across_clues:
            pos = clue.get("position", {})
            row = pos.get("row", 0)
            col = pos.get("col", 0)
            answer = clue.get("answer", "")
            number = clue.get("number", "")

            for i, letter in enumerate(answer):
                if col + i < cols:
                    if grid[row][col + i] is None:
                        grid[row][col + i] = {"letter": "", "number": ""}
                    if i == 0:  # Primera celda lleva el número
                        grid[row][col + i]["number"] = str(number)

        # Marcar las posiciones de las palabras verticales
        for clue in down_clues:
            pos = clue.get("position", {})
            row = pos.get("row", 0)
            col = pos.get("col", 0)
            answer = clue.get("answer", "")
            number = clue.get("number", "")

            for i, letter in enumerate(answer):
                if row + i < rows:
                    if grid[row + i][col] is None:
                        grid[row + i][col] = {"letter": "", "number": ""}
                    if i == 0:  # Primera celda lleva el número
                        if not grid[row + i][col]["number"]:
                            grid[row + i][col]["number"] = str(number)

        # Crear la tabla del crucigrama en Word
        doc.add_heading("Crucigrama", 2)

        # Crear tabla
        table = doc.add_table(rows=rows, cols=cols)
        table.style = 'Table Grid'

        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.rows[row_idx].cells[col_idx]

                if grid[row_idx][col_idx] is None:
                    # Celda negra (no se usa)
                    cell_paragraph = cell.paragraphs[0]
                    cell_paragraph.add_run("")
                    # Fondo oscuro
                    from docx.oxml.ns import qn
                    from docx.oxml import OxmlElement
                    shading_elm = OxmlElement('w:shd')
                    shading_elm.set(qn('w:fill'), '000000')
                    cell._element.get_or_add_tcPr().append(shading_elm)
                else:
                    # Celda blanca con número
                    cell_data = grid[row_idx][col_idx]
                    cell_paragraph = cell.paragraphs[0]

                    if cell_data["number"]:
                        run = cell_paragraph.add_run(cell_data["number"])
                        run.font.size = Pt(8)
                        run.font.bold = True

                    cell_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

        doc.add_paragraph()

        # Pistas Horizontales
        if across_clues:
            doc.add_heading("Pistas Horizontales", 2)
            for clue in across_clues:
                p = doc.add_paragraph()
                p.add_run(f"{clue.get('number', '')}. ").bold = True
                p.add_run(clue.get("clue", ""))
                p.add_run(f" ({len(clue.get('answer', ''))} letras)").italic = True

            doc.add_paragraph()

        # Pistas Verticales
        if down_clues:
            doc.add_heading("Pistas Verticales", 2)
            for clue in down_clues:
                p = doc.add_paragraph()
                p.add_run(f"{clue.get('number', '')}. ").bold = True
                p.add_run(clue.get("clue", ""))
                p.add_run(f" ({len(clue.get('answer', ''))} letras)").italic = True

            doc.add_paragraph()

        # Página nueva para respuestas
        doc.add_page_break()

        # Respuestas (opcional, en una sección separada)
        doc.add_heading("Respuestas", 2)

        if across_clues:
            doc.add_heading("Horizontales:", 3)
            for clue in across_clues:
                doc.add_paragraph(f"{clue.get('number', '')}: {clue.get('answer', '')}", style='List Bullet')

        if down_clues:
            doc.add_heading("Verticales:", 3)
            for clue in down_clues:
                doc.add_paragraph(f"{clue.get('number', '')}: {clue.get('answer', '')}", style='List Bullet')

    def _add_word_search_to_word(self, doc: Document, data: Dict):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except Exception as e:
            print(f"Error parseando contenido de sopa de letras para Word: {e}")
            content = {}

        print(f"Contenido de sopa de letras para Word: {content}")

        # Descripción si existe
        if "description" in content:
            doc.add_paragraph(content["description"])
            doc.add_paragraph()

        # Palabras a buscar
        words = content.get("words", [])
        if words:
            doc.add_heading("Palabras a buscar", 2)
            for word_data in words:
                p = doc.add_paragraph()
                p.add_run(word_data.get("word", "")).bold = True
                if word_data.get("hint"):
                    p.add_run(f" - {word_data.get('hint', '')}")

            doc.add_paragraph()
        else:
            doc.add_heading("Error", 2)
            doc.add_paragraph("No se encontraron palabras para esta sopa de letras.")
            return

        # Cuadrícula (si existe)
        if "grid" in content and content["grid"]:
            doc.add_heading("Sopa de Letras", 2)

            grid = content["grid"]
            # Crear una tabla para la cuadrícula
            table = doc.add_table(rows=len(grid), cols=len(grid[0]) if grid else 0)
            table.style = 'Table Grid'

            for row_idx, grid_row in enumerate(grid):
                for col_idx, letter in enumerate(grid_row):
                    cell = table.rows[row_idx].cells[col_idx]
                    cell.text = letter
                    # Centrar y hacer negrita
                    for paragraph in cell.paragraphs:
                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = Pt(12)

            doc.add_paragraph()
        else:
            doc.add_paragraph("(No se generó la cuadrícula)")

    # Métodos auxiliares para Excel
    def _add_exam_to_excel(self, ws, data: Dict, header_fill, header_font):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        row = 3
        ws[f'A{row}'] = "No."
        ws[f'B{row}'] = "Pregunta"
        ws[f'C{row}'] = "Tipo"
        ws[f'D{row}'] = "Respuesta Correcta"

        for cell in [ws[f'A{row}'], ws[f'B{row}'], ws[f'C{row}'], ws[f'D{row}']]:
            cell.fill = header_fill
            cell.font = header_font

        row += 1
        for q in content.get("questions", []):
            ws[f'A{row}'] = q.get("id", "")
            ws[f'B{row}'] = q.get("question", "")
            ws[f'C{row}'] = q.get("type", "")
            ws[f'D{row}'] = q.get("correct_answer", "")
            row += 1

    def _add_survey_to_excel(self, ws, data: Dict, header_fill, header_font):
        self._add_exam_to_excel(ws, data, header_fill, header_font)

    def _add_rubric_to_excel(self, ws, data: Dict, header_fill, header_font):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except:
            content = {}

        row = 3
        for criterion in content.get("criteria", []):
            ws[f'A{row}'] = criterion.get("name", "")
            ws[f'A{row}'].font = Font(bold=True)
            ws.merge_cells(f'A{row}:D{row}')
            row += 1

            ws[f'A{row}'] = "Nivel"
            ws[f'B{row}'] = "Puntos"
            ws[f'C{row}'] = "Descripción"

            for cell in [ws[f'A{row}'], ws[f'B{row}'], ws[f'C{row}']]:
                cell.fill = header_fill
                cell.font = header_font

            row += 1
            for level in criterion.get("levels", []):
                ws[f'A{row}'] = level.get("level", "")
                ws[f'B{row}'] = level.get("points", "")
                ws[f'C{row}'] = level.get("description", "")
                row += 1

            row += 1

    def _add_crossword_to_excel(self, ws, data: Dict, header_fill, header_font):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except Exception as e:
            print(f"Error parseando contenido del crucigrama: {e}")
            content = {}

        # Debug: Imprimir estructura del contenido
        print(f"Contenido del crucigrama: {content}")

        row = 3
        clues = content.get("clues", {})

        # Verificar si hay pistas
        if not clues:
            ws[f'A{row}'] = "No se encontraron pistas en el crucigrama"
            ws.merge_cells(f'A{row}:D{row}')
            return

        across_clues = clues.get("across", [])
        down_clues = clues.get("down", [])

        if not across_clues and not down_clues:
            ws[f'A{row}'] = "No se generaron pistas para este crucigrama"
            ws.merge_cells(f'A{row}:D{row}')
            return

        # Generar la cuadrícula visual del crucigrama
        grid_size = content.get("grid_size", {"rows": 15, "cols": 15})
        rows_count = grid_size.get("rows", 15)
        cols_count = grid_size.get("cols", 15)

        # Crear matriz de la cuadrícula
        grid = [[None for _ in range(cols_count)] for _ in range(rows_count)]

        # Marcar las posiciones de las palabras horizontales
        for clue in across_clues:
            pos = clue.get("position", {})
            r = pos.get("row", 0)
            c = pos.get("col", 0)
            answer = clue.get("answer", "")
            number = clue.get("number", "")

            for i, letter in enumerate(answer):
                if c + i < cols_count:
                    if grid[r][c + i] is None:
                        grid[r][c + i] = {"letter": "", "number": ""}
                    if i == 0:
                        grid[r][c + i]["number"] = str(number)

        # Marcar las posiciones de las palabras verticales
        for clue in down_clues:
            pos = clue.get("position", {})
            r = pos.get("row", 0)
            c = pos.get("col", 0)
            answer = clue.get("answer", "")
            number = clue.get("number", "")

            for i, letter in enumerate(answer):
                if r + i < rows_count:
                    if grid[r + i][c] is None:
                        grid[r + i][c] = {"letter": "", "number": ""}
                    if i == 0:
                        if not grid[r + i][c]["number"]:
                            grid[r + i][c]["number"] = str(number)

        # Dibujar la cuadrícula en Excel
        ws[f'A{row}'] = "CRUCIGRAMA"
        ws[f'A{row}'].font = Font(size=14, bold=True)
        row += 2

        start_row = row

        for row_idx in range(rows_count):
            for col_idx in range(cols_count):
                excel_row = start_row + row_idx
                excel_col_letter = self._get_excel_column_letter(col_idx)
                cell = ws[f'{excel_col_letter}{excel_row}']

                if grid[row_idx][col_idx] is None:
                    # Celda negra
                    cell.fill = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
                else:
                    # Celda blanca con número
                    cell_data = grid[row_idx][col_idx]
                    cell.value = cell_data.get("number", "")
                    cell.font = Font(size=8, bold=True)
                    cell.fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

                # Añadir bordes
                thin_border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                cell.border = thin_border
                cell.alignment = Alignment(horizontal='left', vertical='top')

                # Ajustar tamaño de celda
                ws.column_dimensions[excel_col_letter].width = 3
                ws.row_dimensions[excel_row].height = 20

        row = start_row + rows_count + 2

        # Pistas Horizontales
        ws[f'A{row}'] = "PISTAS HORIZONTALES"
        ws[f'A{row}'].font = Font(size=12, bold=True)
        row += 1

        for clue in across_clues:
            ws[f'A{row}'] = f"{clue.get('number', '')}. {clue.get('clue', '')} ({len(clue.get('answer', ''))} letras)"
            row += 1

        row += 1

        # Pistas Verticales
        ws[f'A{row}'] = "PISTAS VERTICALES"
        ws[f'A{row}'].font = Font(size=12, bold=True)
        row += 1

        for clue in down_clues:
            ws[f'A{row}'] = f"{clue.get('number', '')}. {clue.get('clue', '')} ({len(clue.get('answer', ''))} letras)"
            row += 1

    def _add_word_search_to_excel(self, ws, data: Dict, header_fill, header_font):
        try:
            content = json.loads(data.get("content", "{}")) if isinstance(data.get("content"), str) else data.get("content", {})
        except Exception as e:
            print(f"Error parseando contenido de sopa de letras: {e}")
            content = {}

        print(f"Contenido de sopa de letras: {content}")

        row = 3

        # Descripción si existe
        if "description" in content:
            ws[f'A{row}'] = "Descripción"
            ws[f'B{row}'] = content["description"]
            ws.merge_cells(f'B{row}:E{row}')
            row += 2

        # Lista de palabras
        ws[f'A{row}'] = "Palabras a buscar"
        ws[f'A{row}'].font = Font(bold=True)
        ws[f'B{row}'] = "Pista"
        ws[f'B{row}'].font = Font(bold=True)
        row += 1

        words = content.get("words", [])
        if words:
            for word_data in words:
                ws[f'A{row}'] = word_data.get("word", "")
                ws[f'B{row}'] = word_data.get("hint", "")
                row += 1
        else:
            ws[f'A{row}'] = "No se encontraron palabras"
            row += 1

        # Agregar la cuadrícula si existe
        if "grid" in content and content["grid"]:
            row += 2
            ws[f'A{row}'] = "Sopa de Letras"
            ws[f'A{row}'].font = Font(bold=True)
            row += 1

            grid = content["grid"]
            for grid_row in grid:
                for col_idx, letter in enumerate(grid_row):
                    cell = ws.cell(row=row, column=col_idx + 1, value=letter)
                    cell.alignment = Alignment(horizontal='center')
                    cell.font = Font(bold=True, size=14)
                    # Añadir bordes a las celdas
                    thin_border = Border(
                        left=Side(style='thin'),
                        right=Side(style='thin'),
                        top=Side(style='thin'),
                        bottom=Side(style='thin')
                    )
                    cell.border = thin_border
                row += 1
        else:
            row += 1
            ws[f'A{row}'] = "No se generó la cuadrícula"


export_service = ExportService()
